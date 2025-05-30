# 文件：backend/services/paper_analyzer/paper_processor.py    实现文档处理
import PyPDF2
from io import BytesIO
import docx
import pptx
import pandas as pd
import markdown
import chardet
from typing import Union, List, Dict, Tuple, Optional, Any
import google.generativeai as genai
from google.genai import types
import io
import os
from dotenv import load_dotenv
from PIL import Image
import base64
from pdf2image import convert_from_bytes
import tempfile
import asyncio
import json
import hashlib
import concurrent.futures
from functools import partial
import re
import subprocess
from .web_processor import WebProcessor, WebProcessingConfig
import logging

logger = logging.getLogger(__name__)

class ProcessingConfig:
    """文档处理配置类 - 增强配置选项"""
    def __init__(self, 
                 extract_for_rag: bool = True,
                 generate_summary: bool = True,
                 max_file_size_mb: int = 50,
                 chunk_size: int = 512,
                 overlap: int = 50,
                 batch_size: int = 5,
                 # 新增配置选项
                 enable_ocr: bool = True,
                 ocr_language: str = 'chi_sim+eng',
                 enable_smart_parsing: bool = True,
                 preserve_formatting: bool = True,
                 extract_tables: bool = True,
                 extract_images: bool = False,
                 quality_threshold: float = 0.8,
                 parallel_processing: bool = True,
                 cache_processing_results: bool = True):
        """
        增强的文档处理配置
        
        新增参数:
        - enable_ocr: 是否启用OCR功能
        - ocr_language: OCR识别语言设置
        - enable_smart_parsing: 启用智能解析（表格、列表等结构）
        - preserve_formatting: 保留原始格式
        - extract_tables: 提取表格数据
        - extract_images: 提取图片（暂不支持）
        - quality_threshold: 处理质量阈值
        - parallel_processing: 并行处理多页文档
        - cache_processing_results: 缓存处理结果
        """
        self.extract_for_rag = extract_for_rag
        self.generate_summary = generate_summary
        self.max_file_size_mb = max_file_size_mb
        self.chunk_size = chunk_size
        self.overlap = overlap
        self.batch_size = batch_size
        
        # 新增配置
        self.enable_ocr = enable_ocr
        self.ocr_language = ocr_language
        self.enable_smart_parsing = enable_smart_parsing
        self.preserve_formatting = preserve_formatting
        self.extract_tables = extract_tables
        self.extract_images = extract_images
        self.quality_threshold = quality_threshold
        self.parallel_processing = parallel_processing
        self.cache_processing_results = cache_processing_results

class PaperProcessor:
    def __init__(self):
        load_dotenv()
        self.api_key = os.getenv("GEMINI_API_KEY")
        if not self.api_key:
            raise ValueError("GEMINI_API_KEY not found in environment variables")
        
        # 配置 Gemini API
        genai.configure(api_key=self.api_key)
        
        # 使用正确的模型名称
        self.model = genai.GenerativeModel('gemini-1.5-pro')
        
        # 配置模型参数
        self.generation_config = {
            "temperature": 0.2,  # 降低创造性以提高准确性
            "top_p": 0.95,
            "top_k": 40,
            "max_output_tokens": 2048,
        }
        
        # 初始化缓存目录
        self.cache_dir = os.path.join(os.path.expanduser("~"), ".cache", "paper_processor")
        os.makedirs(self.cache_dir, exist_ok=True)
        
        # 线程池用于并行处理任务
        self._executor = concurrent.futures.ThreadPoolExecutor(max_workers=4)
        
        # 添加网页处理器
        self.web_processor_config = WebProcessingConfig(
            extract_text=True,
            extract_links=True,
            extract_tables=True,
            extract_code=True,
            chunk_size=512,
            chunk_overlap=50
        )
        
        logger.info(f"[PAPER_PROCESSOR_INIT] 文档处理器初始化完成，包含网页处理能力")

    async def process(self, file_content: bytes, filename: str, extract_for_rag: bool = True) -> dict:
        """处理不同类型的文档文件，返回包含行号信息的内容"""
        try:
            # 配置处理选项
            config = ProcessingConfig(
                extract_for_rag=extract_for_rag,
                generate_summary=True
            )
            
            # 检查文件大小
            file_size_mb = len(file_content) / (1024 * 1024)
            if file_size_mb > config.max_file_size_mb:
                raise ValueError(f"文件过大，超过最大限制 {config.max_file_size_mb}MB")
            
            # 获取文件类型
            extension = filename.lower().split('.')[-1]
            file_type = extension
            
            # 根据文件类型选择处理方法
            structured_data = None  # 初始化结构化数据为None
            
            if extension == 'pdf':
                content, line_mapping = await self._process_pdf(file_content)
            elif extension in ['docx', 'doc']:
                content, line_mapping = await self._process_word(file_content, filename)
            elif extension in ['pptx', 'ppt']:
                content, line_mapping = await self._process_powerpoint(file_content)
            elif extension in ['xlsx', 'xls']:
                # Excel处理会返回额外的结构化数据
                content, line_mapping, structured_data = await self._process_excel(file_content)
            elif extension == 'txt':
                content, line_mapping = await self._process_text(file_content)
            elif extension == 'md':
                content, line_mapping = await self._process_markdown(file_content)
            else:
                raise ValueError(f"Unsupported file type: {extension}")
            
            # 提取是否为扫描件
            is_scanned = any(line.get("is_scanned", False) for line in line_mapping.values())
            
            # 生成文档摘要
            summary = ""
            if config.generate_summary and content:
                summary = await self._generate_document_summary(content, filename)
            
            # 提取关键字/标签
            tags = []
            if config.generate_summary and content:
                tags = await self._extract_document_tags(content, filename)
            
            result = {
                "content": content,
                "line_mapping": line_mapping,
                "total_lines": len(line_mapping),
                "is_scanned": is_scanned,
                "file_type": file_type,
                "file_size": len(file_content),
                "summary": summary,
                "tags": tags
            }
            
            # 如果有结构化数据，添加到结果中
            if structured_data is not None:
                result["structured_data"] = structured_data
            
            return result
                
        except Exception as e:
            raise Exception(f"Paper processing error: {str(e)}")

    async def process_batch(self, files: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """批量处理多个文档
        
        参数:
            files: 文件列表，每个文件为字典：{"content": bytes, "filename": str}
            
        返回:
            处理结果列表
        """
        try:
            # 检查输入
            if not files:
                return []
            
            # 使用半异步方式处理
            results = []
            tasks = []
            
            # 配置批处理大小
            config = ProcessingConfig()
            batch_size = min(config.batch_size, len(files))
            
            # 分批处理
            for i in range(0, len(files), batch_size):
                batch = files[i:i+batch_size]
                batch_tasks = []
                
                for file_info in batch:
                    if not isinstance(file_info, dict) or "content" not in file_info or "filename" not in file_info:
                        continue
                        
                    task = self.process(
                        file_content=file_info["content"],
                        filename=file_info["filename"]
                    )
                    batch_tasks.append(task)
                
                # 并行处理当前批次
                batch_results = await asyncio.gather(*batch_tasks, return_exceptions=True)
                
                # 收集结果
                for j, result in enumerate(batch_results):
                    if isinstance(result, Exception):
                        results.append({
                            "status": "error",
                            "message": str(result),
                            "filename": batch[j]["filename"]
                        })
                    else:
                        result["status"] = "success"
                        result["filename"] = batch[j]["filename"]
                        results.append(result)
            
            return results
            
        except Exception as e:
            print(f"Batch processing error: {str(e)}")
            return [{"status": "error", "message": f"Batch processing error: {str(e)}"}]

    async def _generate_document_summary(self, content: str, filename: str) -> str:
        """生成文档摘要"""
        try:
            # 如果内容太长，截取前后部分
            max_chars = 10000
            if len(content) > max_chars:
                half = max_chars // 2
                summarize_content = content[:half] + "\n\n[...内容过长，中间部分已省略...]\n\n" + content[-half:]
            else:
                summarize_content = content
                
            # 构建提示词
            prompt = f"""请根据以下文档内容，生成一个简洁的摘要（200-300字）。文档文件名为"{filename}"。

文档内容:
{summarize_content}

你的摘要应该:
1. 概括文档的主要内容和目的
2. 保持客观，不添加个人观点
3. 使用简洁清晰的语言
4. 控制在200-300字之间

摘要:"""

            # 调用AI生成摘要
            response = self.model.generate_content(
                prompt,
                generation_config=self.generation_config
            )
            
            if response and response.text:
                return response.text.strip()
            return ""
            
        except Exception as e:
            print(f"Summary generation error: {str(e)}")
            return ""

    async def _extract_document_tags(self, content: str, filename: str) -> List[str]:
        """提取文档关键字/标签"""
        try:
            # 如果内容太长，截取前部分
            max_chars = 8000
            if len(content) > max_chars:
                content_to_analyze = content[:max_chars]
            else:
                content_to_analyze = content
                
            # 构建提示词
            prompt = f"""请根据以下文档内容，提取5-10个关键词或标签。文档文件名为"{filename}"。

文档内容:
{content_to_analyze}

要求:
1. 提取能够代表文档主题的关键词或短语
2. 每个标签不超过3个词
3. 以英文逗号分隔列出
4. 不要编号或使用其他格式

关键词:"""

            # 调用AI生成标签
            response = self.model.generate_content(
                prompt,
                generation_config={
                    "temperature": 0.1,
                    "max_output_tokens": 256,
                }
            )
            
            if response and response.text:
                # 清理并分割标签
                tags_text = response.text.strip()
                
                # 移除可能的引号、编号、破折号等
                tags_text = re.sub(r'^[\'"]*|[\'"]*$', '', tags_text)
                tags_text = re.sub(r'^\d+\.\s*', '', tags_text, flags=re.MULTILINE)
                tags_text = re.sub(r'^-\s*', '', tags_text, flags=re.MULTILINE)
                
                # 分割并清理标签
                tags = [tag.strip() for tag in re.split(r',|\n', tags_text) if tag.strip()]
                
                # 限制标签数量
                return tags[:10]
            return []
            
        except Exception as e:
            print(f"Tags extraction error: {str(e)}")
            return []

    async def _process_pdf(self, file_content: bytes) -> tuple:
        """处理 PDF 文件，支持扫描件"""
        try:
            # 1. 首先尝试使用 PyPDF2 提取文本
            pdf_file = BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            # 检查是否为扫描件
            is_scanned = False
            for page in pdf_reader.pages:
                if not page.extract_text().strip():
                    is_scanned = True
                    break
            
            if is_scanned:
                # 2. 如果是扫描件，使用 Gemini API 处理
                return await self._process_scanned_pdf(file_content)
            else:
                # 3. 如果是普通 PDF，使用原有逻辑
                return await self._process_normal_pdf(pdf_reader)
                
        except Exception as e:
            raise Exception(f"PDF processing error: {str(e)}")

    async def _process_scanned_pdf(self, file_content: bytes) -> tuple:
        """使用 Gemini API 处理扫描件 PDF"""
        try:
            # 1. 将 PDF 转换为图片
            images = convert_from_bytes(file_content)
            
            # 2. 处理每一页
            all_content = []
            line_mapping = {}
            current_line = 1
            
            # 创建任务列表
            tasks = []
            for page_num, image in enumerate(images, 1):
                # 将图片转换为字节数据
                img_byte_arr = io.BytesIO()
                image.save(img_byte_arr, format='JPEG')
                img_byte_arr = img_byte_arr.getvalue()
                
                # 创建处理任务
                task = self._process_scanned_pdf_page(img_byte_arr, page_num)
                tasks.append(task)
            
            # 并行处理所有页面
            results = await asyncio.gather(*tasks)
            
            # 组合结果
            for page_num, page_content in results:
                if page_content:
                    all_content.append(page_content)
                    
                    # 处理行号映射
                    lines = page_content.split('\n')
                    for line in lines:
                        if line.strip():
                            line_mapping[current_line] = {
                                "content": line,
                                "page": page_num,
                                "start_pos": len('\n'.join(all_content)) - len(line) - 1,
                                "end_pos": len('\n'.join(all_content)) - 1,
                                "is_scanned": True
                            }
                            current_line += 1
            
            content = '\n'.join(all_content)
            return content, line_mapping
            
        except Exception as e:
            raise Exception(f"Scanned PDF processing error: {str(e)}")

    async def _process_scanned_pdf_page(self, img_data: bytes, page_num: int) -> Tuple[int, str]:
        """处理单个扫描PDF页面"""
        try:
            # 生成缓存键
            cache_key = hashlib.md5(img_data).hexdigest()
            cache_file = os.path.join(self.cache_dir, f"page_{cache_key}.txt")
            
            # 检查缓存
            if os.path.exists(cache_file):
                with open(cache_file, 'r', encoding='utf-8') as f:
                    return page_num, f.read()
            
            # 调用 Gemini API 处理图片
            doc_data = {
                "mime_type": "image/jpeg",
                "data": img_data
            }
            
            prompt = f"""请提取这个图片中的文本内容（第{page_num}页），包括：
            1. 正文内容
            2. 表格内容
            3. 图片中的文字
            4. 页眉页脚
            请保持原文的段落结构和格式。"""
            
            response = self.model.generate_content(
                contents=[doc_data, prompt]
            )
            
            page_content = ""
            if response and response.text:
                page_content = response.text
                
                # 保存到缓存
                with open(cache_file, 'w', encoding='utf-8') as f:
                    f.write(page_content)
            
            return page_num, page_content
            
        except Exception as e:
            print(f"Page {page_num} processing error: {str(e)}")
            return page_num, ""

    async def _process_normal_pdf(self, pdf_reader: PyPDF2.PdfReader) -> tuple:
        """处理普通 PDF 文件（原有逻辑）"""
        content = ""
        line_mapping = {}
        current_line = 1
        
        for page_num, page in enumerate(pdf_reader.pages, 1):
            page_text = page.extract_text()
            lines = page_text.split('\n')
            
            for line in lines:
                if line.strip():
                    content += line + "\n"
                    line_mapping[current_line] = {
                        "content": line,
                        "page": page_num,
                        "start_pos": len(content) - len(line) - 1,
                        "end_pos": len(content) - 1,
                        "is_scanned": False  # 标记为非扫描件
                    }
                    current_line += 1
                    
        return content, line_mapping

    async def _process_word(self, file_content: bytes, filename: str) -> tuple:
        """处理 Word 文件，返回内容和行号映射"""
        try:
            # 获取扩展名
            extension = filename.lower().split('.')[-1]
            
            # 处理.docx格式
            if extension == 'docx':
                doc_file = BytesIO(file_content)
                doc = docx.Document(doc_file)
                
                content = ""
                line_mapping = {}
                current_line = 1
                
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        content += paragraph.text + "\n"
                        line_mapping[current_line] = {
                            "content": paragraph.text,
                            "page": 1,
                            "start_pos": len(content) - len(paragraph.text) - 1,
                            "end_pos": len(content) - 1,
                            "is_scanned": False  # 添加这个字段
                        }
                        current_line += 1
                
                # 处理表格
                for table in doc.tables:
                    for row in table.rows:
                        row_text = " | ".join(cell.text for cell in row.cells)
                        if row_text.strip():
                            content += row_text + "\n"
                            line_mapping[current_line] = {
                                "content": row_text,
                                "page": 1,
                                "start_pos": len(content) - len(row_text) - 1,
                                "end_pos": len(content) - 1,
                                "is_scanned": False  # 添加这个字段
                            }
                            current_line += 1
                
                return content, line_mapping
            
            # 处理.doc格式（旧版Word）
            elif extension == 'doc':
                # 考虑使用其他库处理.doc文件
                # 例如：使用textract或其他兼容工具
                # 临时解决方案：使用通用文本提取模式
                return await self._extract_doc_content(file_content)
            
        except Exception as e:
            raise Exception(f"Word processing error: {str(e)}")

    async def _extract_doc_content(self, file_content: bytes) -> tuple:
        """尝试从旧版Word文档(.doc)中提取文本"""
        try:
            # 这里需要添加依赖：pip install textract 或其他替代方案
            import tempfile
            import subprocess
            
            # 创建临时文件
            with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as temp:
                temp_path = temp.name
                temp.write(file_content)
            
            # 使用外部工具antiword提取文本（需安装）
            try:
                content = subprocess.check_output(['antiword', temp_path]).decode('utf-8')
            except:
                # 降级方案：尝试使用strings提取一些文本
                content = subprocess.check_output(['strings', temp_path]).decode('utf-8')
                
            # 清理临时文件
            os.unlink(temp_path)
            
            # 构建简单行映射
            lines = content.split('\n')
            line_mapping = {}
            current_pos = 0
            
            for i, line in enumerate(lines, 1):
                content_len = len(line)
                if line.strip():
                    line_mapping[i] = {
                        "content": line,
                        "page": 1,
                        "start_pos": current_pos,
                        "end_pos": current_pos + content_len,
                        "is_scanned": False
                    }
                current_pos += content_len + 1  # +1 for newline
            
            return content, line_mapping
            
        except Exception as e:
            raise Exception(f"DOC extraction error: {str(e)}")

    async def _process_powerpoint(self, file_content: bytes) -> tuple:
        """处理 PowerPoint 文件，返回内容和行号映射"""
        try:
            ppt_file = BytesIO(file_content)
            ppt = pptx.Presentation(ppt_file)
            
            content = ""
            line_mapping = {}
            current_line = 1
            
            # 添加幻灯片页码
            slide_count = len(ppt.slides)
            
            for slide_num, slide in enumerate(ppt.slides, 1):
                # 添加幻灯片页码分隔符
                content += f"\n[Slide {slide_num}/{slide_count}]\n"
                
                # 处理幻灯片标题
                if slide.shapes.title:
                    title_text = f"Title: {slide.shapes.title.text}"
                    content += title_text + "\n"
                    line_mapping[current_line] = {
                        "content": title_text,
                        "page": slide_num,
                        "start_pos": len(content) - len(title_text) - 1,
                        "end_pos": len(content) - 1,
                        "is_scanned": False
                    }
                    current_line += 1
                
                # 处理幻灯片内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        shape_text = shape.text.strip()
                        if shape_text:
                            content += shape_text + "\n"
                            lines = shape_text.split('\n')
                            
                            for line in lines:
                                if line.strip():
                                    line_mapping[current_line] = {
                                        "content": line,
                                        "page": slide_num,
                                        "start_pos": len(content) - len(line) - 1,
                                        "end_pos": len(content) - 1,
                                        "is_scanned": False
                                    }
                                    current_line += 1
                
            return content, line_mapping
            
        except Exception as e:
            raise Exception(f"PowerPoint processing error: {str(e)}")

    async def _process_excel(self, file_content: bytes) -> tuple:
        """处理 Excel 文件，返回内容和行号映射，并存储结构化数据"""
        try:
            excel_file = BytesIO(file_content)
            
            # 读取所有工作表
            content = ""
            line_mapping = {}
            current_line = 1
            
            # 存储结构化数据
            structured_data = {}
            
            try:
                excel = pd.ExcelFile(excel_file)
                
                for sheet_name in excel.sheet_names:
                    try:
                        # 尝试多种解析方式，优化Excel读取
                        try:
                            # 先尝试标准读取（有可能有表头）
                            df = pd.read_excel(
                                excel, 
                                sheet_name=sheet_name,
                                dtype=str,  # 避免数值被转为科学计数法
                                na_filter=False,  # 避免空值被替换为NaN
                            )
                            
                            # 检查是否有可能的表头问题
                            header_issue = False
                            
                            # 检查所有列名是否都是数字（可能是误判的表头）
                            if df.shape[0] > 0 and all(isinstance(col, (int, float)) for col in df.columns):
                                header_issue = True
                                
                            # 检查第一行是否像表头（与列名相似）
                            if df.shape[0] > 0 and not header_issue:
                                first_row = df.iloc[0].astype(str).tolist()
                                column_names = [str(col) for col in df.columns]
                                similarity = sum(1 for a, b in zip(first_row, column_names) if a == b)
                                if similarity / len(column_names) > 0.5:  # 如果超过50%相似
                                    header_issue = True
                            
                            # 如果检测到表头可能有问题，重新读取
                            if header_issue:
                                df = pd.read_excel(
                                    excel,
                                    sheet_name=sheet_name,
                                    header=None,  # 不使用第一行作为表头
                                    dtype=str,
                                    na_filter=False
                                )
                        except Exception:
                            # 如果常规方法失败，尝试更保守的方法
                            df = pd.read_excel(
                                excel,
                                sheet_name=sheet_name,
                                header=None,
                                dtype=str,
                                na_filter=False
                            )
                        
                        # 处理空表格的情况
                        if df.empty:
                            sheet_header = f"\nSheet: {sheet_name} [空表格]\n"
                            content += sheet_header
                            
                            line_mapping[current_line] = {
                                "content": sheet_header.strip(),
                                "page": 1,
                                "start_pos": len(content) - len(sheet_header),
                                "end_pos": len(content),
                                "is_scanned": False
                            }
                            current_line += 1
                            structured_data[sheet_name] = []
                            # 添加空表元数据
                            structured_data[f"{sheet_name}_metadata"] = {
                                "total_rows": 0,
                                "total_columns": 0,
                                "data_range": "0-0",
                                "column_types": {}
                            }
                            continue
                        
                        # 数据清洗和规范化
                        # 1. 移除完全空的行和列
                        df = df.dropna(how='all').dropna(axis=1, how='all')
                        
                        # 2. 如果某列都是NaN或空字符串，填充为空字符串以避免丢失列结构
                        df = df.fillna('')
                        
                        # 3. 处理非字符串列
                        for col in df.columns:
                            # 处理日期时间列，确保格式一致
                            if pd.api.types.is_datetime64_any_dtype(df[col]):
                                df[col] = df[col].dt.strftime('%Y-%m-%d %H:%M:%S')
                            # 确保不存在科学计数法
                            elif pd.api.types.is_numeric_dtype(df[col]):
                                df[col] = df[col].astype(str).apply(
                                    lambda x: x.replace('.0', '') if x.endswith('.0') else x
                                )
                        
                        # 创建列数据类型和统计信息字典
                        column_metadata = {}
                        for col in df.columns:
                            column_values = df[col].astype(str).tolist()
                            non_empty_values = [v for v in column_values if v and v.strip()]
                            
                            # 确定列数据类型
                            if all(v.replace('.', '', 1).isdigit() for v in non_empty_values if v != ''):
                                col_type = "numeric"
                                # 尝试将值转换为数字以进行统计
                                try:
                                    numeric_values = [float(v) for v in non_empty_values if v != '']
                                    if numeric_values:
                                        column_metadata[col] = {
                                            "type": col_type,
                                            "count": len(numeric_values),
                                            "min": min(numeric_values),
                                            "max": max(numeric_values),
                                            "avg": sum(numeric_values) / len(numeric_values),
                                            "non_empty_count": len(non_empty_values),
                                            "empty_count": len(column_values) - len(non_empty_values)
                                        }
                                    else:
                                        column_metadata[col] = {"type": col_type, "count": 0}
                                except:
                                    column_metadata[col] = {"type": "text", "count": len(non_empty_values)}
                            elif any(re.match(r'\d{4}[-/]\d{1,2}[-/]\d{1,2}', v) for v in non_empty_values):
                                column_metadata[col] = {
                                    "type": "date",
                                    "count": len(non_empty_values),
                                    "non_empty_count": len(non_empty_values),
                                    "empty_count": len(column_values) - len(non_empty_values)
                                }
                            else:
                                column_metadata[col] = {
                                    "type": "text",
                                    "count": len(non_empty_values),
                                    "non_empty_count": len(non_empty_values), 
                                    "empty_count": len(column_values) - len(non_empty_values),
                                    "unique_count": len(set(non_empty_values))
                                }
                        
                        # 储存结构化数据 (orient="records"格式)
                        sheet_data = df.to_dict(orient="records")
                        structured_data[sheet_name] = sheet_data
                        
                        # 添加元数据
                        structured_data[f"{sheet_name}_metadata"] = {
                            "total_rows": len(df),
                            "total_columns": len(df.columns),
                            "data_range": f"1-{len(df)}",
                            "columns": list(df.columns),
                            "column_metadata": column_metadata
                        }
                        
                        # 添加工作表标题信息
                        rows_count = len(df)
                        cols_count = len(df.columns)
                        sheet_header = f"\nSheet: {sheet_name} [行数: {rows_count}, 列数: {cols_count}]\n"
                        content += sheet_header
                        
                        # 添加工作表标题行
                        line_mapping[current_line] = {
                            "content": sheet_header.strip(),
                            "page": 1,
                            "start_pos": len(content) - len(sheet_header),
                            "end_pos": len(content),
                            "is_scanned": False
                        }
                        current_line += 1
                        
                        # 转换表格内容为文本，确保完整展示
                        table_str = df.to_string(
                            index=False, 
                            col_space=30,  # 确保字段完整展示
                            max_rows=None,  # 不限制行数
                            max_cols=None,  # 不限制列数
                            min_rows=df.shape[0]  # 确保所有行都显示
                        )
                        content += table_str + "\n\n"
                        lines = table_str.split('\n')
                        
                        for line in lines:
                            if line.strip():
                                line_mapping[current_line] = {
                                    "content": line,
                                    "page": 1,
                                    "start_pos": len(content) - len(line) - 2,
                                    "end_pos": len(content) - 2,
                                    "is_scanned": False,
                                    "sheet_name": sheet_name  # 添加工作表名称
                                }
                                current_line += 1
                    
                    except Exception as sheet_error:
                        print(f"Error processing sheet {sheet_name}: {str(sheet_error)}")
                        # 记录错误信息，但继续处理其他表格
                        error_msg = f"\nSheet: {sheet_name} [处理错误: {str(sheet_error)}]\n"
                        content += error_msg
                        line_mapping[current_line] = {
                            "content": error_msg.strip(),
                            "page": 1,
                            "start_pos": len(content) - len(error_msg),
                            "end_pos": len(content),
                            "is_scanned": False
                        }
                        current_line += 1
                        structured_data[sheet_name] = {"error": str(sheet_error)}
                
            except Exception as excel_error:
                # 处理Excel读取失败的情况
                error_msg = f"Excel文件读取失败: {str(excel_error)}"
                content += error_msg
                line_mapping[current_line] = {
                    "content": error_msg,
                    "page": 1,
                    "start_pos": len(content) - len(error_msg),
                    "end_pos": len(content),
                    "is_scanned": False
                }
                current_line += 1
            
            # 返回包含结构化数据的结果
            return content, line_mapping, structured_data
            
        except Exception as e:
            raise Exception(f"Excel processing error: {str(e)}")

    async def _process_text(self, file_content: bytes) -> tuple:
        """处理文本文件，返回内容和行号映射"""
        try:
            # 尝试检测编码
            encoding = chardet.detect(file_content)['encoding'] or 'utf-8'
            
            try:
                content = file_content.decode(encoding)
            except UnicodeDecodeError:
                # 如果检测到的编码不正确，尝试常用编码
                for enc in ['utf-8', 'latin-1', 'gbk', 'cp1252', 'iso-8859-1']:
                    try:
                        content = file_content.decode(enc)
                        break
                    except UnicodeDecodeError:
                        continue
                else:
                    # 如果所有尝试都失败，使用latin-1（不会抛出解码错误）
                    content = file_content.decode('latin-1')
            
            lines = content.split('\n')
            
            line_mapping = {}
            current_pos = 0
            
            for i, line in enumerate(lines, 1):
                content_len = len(line)
                if line.strip():
                    line_mapping[i] = {
                        "content": line,
                        "page": 1,
                        "start_pos": current_pos,
                        "end_pos": current_pos + content_len,
                        "is_scanned": False
                    }
                current_pos += content_len + 1  # +1 for newline
            
            return content, line_mapping
            
        except Exception as e:
            raise Exception(f"Text processing error: {str(e)}")

    async def _process_markdown(self, file_content: bytes) -> tuple:
        """处理 Markdown 文件，返回内容和行号映射"""
        try:
            # 首先获取文本内容
            content, line_mapping = await self._process_text(file_content)
            
            # 将 Markdown 转换为纯文本
            html = markdown.markdown(content)
            
            # HTML 到纯文本的转换逻辑
            plain_text = self._html_to_text(html)
            
            # 创建新的行号映射
            new_mapping = {}
            current_line = 1
            
            for line in plain_text.split('\n'):
                if line.strip():
                    new_mapping[current_line] = {
                        "content": line,
                        "page": 1,
                        "start_pos": plain_text.find(line),
                        "end_pos": plain_text.find(line) + len(line),
                        "is_scanned": False
                    }
                    current_line += 1
            
            return plain_text, new_mapping
            
        except Exception as e:
            raise Exception(f"Markdown processing error: {str(e)}")

    def _html_to_text(self, html: str) -> str:
        """将HTML转换为纯文本，保留基本格式"""
        # 移除HTML标签但保留内容
        text = re.sub(r'<head[^>]*>.*?</head>', '', html, flags=re.DOTALL)
        
        # 在<br>标签处添加换行
        text = re.sub(r'<br[^>]*>', '\n', text)
        
        # 在段落标签结束处添加两个换行
        text = re.sub(r'</p>', '\n\n', text)
        
        # 在标题标签结束处添加换行
        text = re.sub(r'</h[1-6]>', '\n\n', text)
        
        # 在列表项标签结束处添加换行
        text = re.sub(r'</li>', '\n', text)
        
        # 移除所有剩余的HTML标签
        text = re.sub(r'<[^>]+>', '', text)
        
        # 处理HTML实体
        text = text.replace('&nbsp;', ' ')
        text = text.replace('&lt;', '<')
        text = text.replace('&gt;', '>')
        text = text.replace('&amp;', '&')
        text = text.replace('&quot;', '"')
        
        # 移除多余的空行
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        return text.strip()

    async def process_url(self, url: str, **kwargs) -> dict:
        """处理网页URL - 新增方法"""
        try:
            logger.info(f"[PROCESS_URL] 开始处理URL: {url}")
            
            async with WebProcessor(self.web_processor_config) as web_processor:
                # 处理网页内容
                web_result = await web_processor.process_url(url)
                
                # 转换为统一格式
                processed_result = {
                    "content": web_result["content"],
                    "line_mapping": self._convert_web_chunks_to_line_mapping(web_result["chunks"]),
                    "total_lines": len(web_result["chunks"]),
                    "structured_data": web_result["structured_data"],
                    "metadata": web_result["metadata"],
                    "source_url": url,
                    "title": web_result["title"]
                }
                
                logger.info(f"[PROCESS_URL_SUCCESS] URL处理成功: {url}")
                return processed_result
                
        except Exception as e:
            logger.error(f"[PROCESS_URL_ERROR] URL处理失败: {url} - {str(e)}")
            raise Exception(f"URL处理失败: {str(e)}")

    async def process_multiple_urls(self, urls: List[str], **kwargs) -> List[dict]:
        """批量处理多个URL - 新增方法"""
        try:
            logger.info(f"[PROCESS_URLS_BATCH] 开始批量处理URL - 数量: {len(urls)}")
            
            async with WebProcessor(self.web_processor_config) as web_processor:
                # 批量处理
                web_results = await web_processor.process_multiple_urls(urls)
                
                # 转换为统一格式
                processed_results = []
                for web_result in web_results:
                    if not web_result.get("error"):
                        processed_result = {
                            "content": web_result["content"],
                            "line_mapping": self._convert_web_chunks_to_line_mapping(web_result["chunks"]),
                            "total_lines": len(web_result["chunks"]),
                            "structured_data": web_result["structured_data"],
                            "metadata": web_result["metadata"],
                            "source_url": web_result["url"],
                            "title": web_result["title"]
                        }
                        processed_results.append(processed_result)
                    else:
                        logger.error(f"[PROCESS_URLS_BATCH_ERROR] URL处理失败: {web_result['url']}")
                
                logger.info(f"[PROCESS_URLS_BATCH_SUCCESS] 批量URL处理完成 - 成功: {len(processed_results)}/{len(urls)}")
                return processed_results
                
        except Exception as e:
            logger.error(f"[PROCESS_URLS_BATCH_ERROR] 批量URL处理失败: {str(e)}")
            raise Exception(f"批量URL处理失败: {str(e)}")

    def _convert_web_chunks_to_line_mapping(self, chunks: List[Dict]) -> Dict[str, Dict]:
        """将网页分块转换为line_mapping格式"""
        line_mapping = {}
        
        for i, chunk in enumerate(chunks):
            line_key = str(i + 1)  # 行号从1开始
            line_mapping[line_key] = {
                "content": chunk.get("content", ""),
                "start_pos": chunk.get("start_pos", 0),
                "end_pos": chunk.get("end_pos", 0),
                "chunk_index": chunk.get("chunk_index", i),
                "source_url": chunk.get("source_url", ""),
                "is_web_chunk": True  # 标记为网页分块
            }
        
        return line_mapping
